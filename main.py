api='AIzaSyCWuqsvu0pmNSgUgRJIRYCl61TzZDVdWMs'


import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from docx import Document
from io import BytesIO

genai.configure(api_key=api)


def generate_content(prompt, instruction=""):
    model = genai.GenerativeModel("gemini-1.5-pro-002")
    if instruction:
        full_prompt = f"{prompt}\nInstruction: {instruction}"
    else:
        full_prompt = prompt
    response = model.generate_content(full_prompt)
    return response.text


def create_pptx(content, heading):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = heading

    textbox = slide.shapes.placeholders[1]
    textbox.text = content

    pptx_file = BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)
    return pptx_file


def create_ppt(content, heading):
    return create_pptx(content, heading)


def create_docx(content, heading):
    doc = Document()
    doc.add_heading(heading, 0)
    doc.add_paragraph(content)

    docx_file = BytesIO()
    doc.save(docx_file)
    docx_file.seek(0)
    return docx_file


def app():
    st.title("Assignment Generator using AI")

    if 'content' not in st.session_state:
        st.session_state.content = None

    if 'prompt' not in st.session_state:
        st.session_state.prompt = ''
    if 'instruction' not in st.session_state:
        st.session_state.instruction = ''
    if 'export_format' not in st.session_state:
        st.session_state.export_format = 'PPTX'
    if 'heading' not in st.session_state:
        st.session_state.heading = "Generated Content"

    st.session_state.prompt = st.text_area("Enter the prompt:", value=st.session_state.prompt, height=150)
    st.session_state.instruction = st.text_area("Enter additional instructions (optional):", value=st.session_state.instruction, height=150)

    if st.button("Generate"):
        if st.session_state.prompt:
            with st.spinner("Generating content..."):
                st.session_state.content = generate_content(st.session_state.prompt, st.session_state.instruction)
                st.session_state.heading = st.session_state.prompt[:30] + "..." if len(st.session_state.prompt) > 30 else st.session_state.prompt
                st.success("Content generated successfully!")

    if st.session_state.content:
        st.session_state.export_format = st.radio("Select export format:", ("PPTX", "PPT", "DOCX"), index=["PPTX", "PPT", "DOCX"].index(st.session_state.export_format))

        download_filename = st.session_state.prompt[:30] + "..." if len(st.session_state.prompt) > 30 else st.session_state.prompt
        download_filename = download_filename.replace(" ", "_")  # Replace spaces with underscores to avoid issues with filenames
        download_filename = download_filename[:50]  # Limit filename length to avoid issues

        if st.session_state.export_format == "PPTX":
            pptx_file = create_pptx(st.session_state.content, st.session_state.heading)
            st.download_button(
                label="Download PPTX",
                data=pptx_file,
                file_name=f"{download_filename}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        elif st.session_state.export_format == "PPT":
            ppt_file = create_ppt(st.session_state.content, st.session_state.heading)  # This will generate a PPTX but could be manually converted
            st.download_button(
                label="Download PPT (PPTX format, please convert manually if needed)",
                data=ppt_file,
                file_name=f"{download_filename}.ppt",
                mime="application/vnd.ms-powerpoint"
            )
        elif st.session_state.export_format == "DOCX":
            docx_file = create_docx(st.session_state.content, st.session_state.heading)
            st.download_button(
                label="Download DOCX",
                data=docx_file,
                file_name=f"{download_filename}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )


if __name__ == "__main__":
    app()

